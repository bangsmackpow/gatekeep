import logging
import hashlib
from pathlib import Path
from datetime import datetime
from typing import Optional
from src.ingestion.models import ExtractedDocument

logger = logging.getLogger(__name__)


def parse_docx(file_path: str | bytes, filename: str = "document.docx") -> ExtractedDocument:
    from docx import Document
    import io

    if isinstance(file_path, bytes):
        doc = Document(io.BytesIO(file_path))
        content = file_path
    else:
        doc = Document(file_path)
        with open(file_path, "rb") as f:
            content = f.read()

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    extracted_text = "\n".join(paragraphs)

    core_properties = doc.core_properties
    created = core_properties.created
    modified = core_properties.modified

    return ExtractedDocument(
        original_filename=filename,
        file_extension="docx",
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        file_size_bytes=len(content),
        sha256_hash=hashlib.sha256(content).hexdigest(),
        content_bytes=b"",
        author=core_properties.author,
        title=core_properties.title,
        subject=core_properties.subject,
        created_date=created if isinstance(created, datetime) else None,
        modified_date=modified if isinstance(modified, datetime) else None,
        extracted_text=extracted_text,
    )


def parse_xlsx(file_path: str | bytes, filename: str = "spreadsheet.xlsx") -> ExtractedDocument:
    import openpyxl
    import io

    if isinstance(file_path, bytes):
        wb = openpyxl.load_workbook(io.BytesIO(file_path), read_only=True, data_only=True)
        content = file_path
    else:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        with open(file_path, "rb") as f:
            content = f.read()

    text_parts = []
    for sheet in wb.worksheets:
        text_parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text_parts.append(row_text)

    extracted_text = "\n".join(text_parts)
    wb.close()

    return ExtractedDocument(
        original_filename=filename,
        file_extension="xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        file_size_bytes=len(content),
        sha256_hash=hashlib.sha256(content).hexdigest(),
        content_bytes=b"",
        author=wb.properties.creator,
        title=wb.properties.title,
        extracted_text=extracted_text,
    )


def parse_pptx(file_path: str | bytes, filename: str = "presentation.pptx") -> ExtractedDocument:
    from pptx import Presentation
    import io

    if isinstance(file_path, bytes):
        prs = Presentation(io.BytesIO(file_path))
        content = file_path
    else:
        prs = Presentation(file_path)
        with open(file_path, "rb") as f:
            content = f.read()

    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text_parts.append(paragraph.text)

    core_properties = prs.core_properties

    return ExtractedDocument(
        original_filename=filename,
        file_extension="pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        file_size_bytes=len(content),
        sha256_hash=hashlib.sha256(content).hexdigest(),
        content_bytes=b"",
        author=core_properties.author,
        title=core_properties.title,
        subject=core_properties.subject,
        created_date=core_properties.created if isinstance(core_properties.created, datetime) else None,
        modified_date=core_properties.modified if isinstance(core_properties.modified, datetime) else None,
        extracted_text="\n".join(text_parts),
    )


def parse_office_file(file_path: str | bytes, filename: str) -> ExtractedDocument:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    parsers = {
        "docx": parse_docx,
        "xlsx": parse_xlsx,
        "pptx": parse_pptx,
    }

    parser = parsers.get(ext)
    if parser:
        return parser(file_path, filename)

    raise ValueError(f"Unsupported Office format: .{ext}")
